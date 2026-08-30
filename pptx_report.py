"""Regenerates the monthly review presentation from its existing template,
swapping only the chart data (categories + series values) for the current
month's numbers. Layout, colors, titles, and every non-chart slide are left
untouched — those slides (capacity upgrades, relocations, new sites, Q&A)
are edited by hand each month, same as before.

Data source is the already-generated Excel report (build_excel_report's
output), not the raw fetch, so the presentation always matches exactly what
is in the accompanying spreadsheet.
"""
import re
from datetime import date

import pandas as pd
from pptx import Presentation
from pptx.oxml.ns import qn

EXCEL_EPOCH = date(1899, 12, 30)


def _to_excel_serial(d):
    return (d - EXCEL_EPOCH).days

# (substring checked in this order; first match wins) -> Excel sheet name
TITLE_TO_SHEET = [
    ("total voice", "Total Voice"),
    ("total data", "Total Data"),
    ("2g voice", "2G Voice"),
    ("3g voice", "3G Voice"),
    ("4g data", "4G Data"),
    ("3g data", "3G Data"),
    ("volte", None),  # resolved below: "network" in title -> Network sheet, else Cluster sheet
]


def _resolve_sheet_name(title, available_sheets):
    title_lower = title.lower()
    for keyword, sheet_name in TITLE_TO_SHEET:
        if keyword not in title_lower:
            continue
        if keyword == "volte":
            sheet_name = "4G Voice - Network" if "network" in title_lower else "4G Voice - Cluster"
        if sheet_name in available_sheets:
            return sheet_name
        return None
    return None


def _load_sheets(excel_path):
    xl = pd.ExcelFile(excel_path)
    sheets = {}
    for name in xl.sheet_names:
        df = xl.parse(name)
        df["Date"] = pd.to_datetime(df["Date"]).dt.date
        df.columns = [c.strip() if isinstance(c, str) else c for c in df.columns]
        sheets[name] = df
    return sheets


def _chart_title(chart):
    try:
        return chart.chart_title.text_frame.text if chart.has_title else ""
    except Exception:
        return ""


def _find(el, *tags):
    return el.find("/".join(qn(t) for t in tags))


def _rewrite_num_cache(num_cache_el, values, is_date=False):
    for child in num_cache_el.findall(qn("c:pt")):
        num_cache_el.remove(child)
    existing_count = num_cache_el.find(qn("c:ptCount"))
    if existing_count is not None:
        num_cache_el.remove(existing_count)

    ptCount_el = num_cache_el.makeelement(qn("c:ptCount"), {"val": str(len(values))})
    num_cache_el.append(ptCount_el)
    for idx, v in enumerate(values):
        v = _to_excel_serial(v) if is_date else v
        pt_el = num_cache_el.makeelement(qn("c:pt"), {"idx": str(idx)})
        v_el = pt_el.makeelement(qn("c:v"), {})
        v_el.text = str(v)
        pt_el.append(v_el)
        num_cache_el.append(pt_el)


def _bump_formula_range(ref_el, new_count):
    f_el = ref_el.find(qn("c:f"))
    if f_el is None or not f_el.text:
        return
    m = re.match(r"^(.*\$[A-Z]+\$)\d+(\$?)$", f_el.text)
    if m:
        f_el.text = f"{m.group(1)}{1 + new_count}{m.group(2)}"


def update_chart_data(chart, df, log):
    dates = list(df["Date"])

    series_cols = {}
    missing = []
    for series in chart.series:
        col = series.name.strip() if series.name else series.name
        if col not in df.columns:
            missing.append(series.name)
        else:
            series_cols[series] = col

    if missing:
        raise ValueError(f"series {missing} not found in sheet columns {list(df.columns)}")

    for series, col in series_cols.items():
        ser_el = series._element
        cat_numRef = _find(ser_el, "c:cat", "c:numRef")
        val_numRef = _find(ser_el, "c:val", "c:numRef")
        cat_cache = _find(ser_el, "c:cat", "c:numRef", "c:numCache")
        val_cache = _find(ser_el, "c:val", "c:numRef", "c:numCache")

        _rewrite_num_cache(cat_cache, dates, is_date=True)
        _rewrite_num_cache(val_cache, df[col].fillna(0).tolist())
        _bump_formula_range(cat_numRef, len(dates))
        _bump_formula_range(val_numRef, len(dates))

    log(f"    updated ({len(chart.series)} series x {len(dates)} dates)")


def build_pptx_report(template_path, excel_path, output_path, log=print):
    sheets = _load_sheets(excel_path)
    prs = Presentation(template_path)

    updated, skipped = 0, 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            chart = shape.chart
            title = _chart_title(chart)
            sheet_name = _resolve_sheet_name(title, sheets)
            if sheet_name is None:
                log(f"slide {i+1}: chart {title!r} -> no matching sheet, left as-is")
                skipped += 1
                continue
            log(f"slide {i+1}: chart {title!r} -> {sheet_name!r}")
            update_chart_data(chart, sheets[sheet_name], log)
            updated += 1

    prs.save(output_path)
    log(f"Done. {updated} charts updated, {skipped} left untouched. Saved to {output_path}")
    return output_path
